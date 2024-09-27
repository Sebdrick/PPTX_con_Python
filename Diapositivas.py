from pptx import Presentation

file_path = 'Presentación.pptx'

wildcard_replacement = {
    "%title%": "Las mejores sagas de terror",
    "%nombre%": "Sebastián Solis Villafuerte",
    "%subtitle1%": "Halloween (1978 – 2022)",
    "%subtitle2%": "Viernes 13 (1980 – 2009)",
    "%subtitle3%": "Pesadilla en Elm Street (1980 – 2010)",
    "%subtitle4%": "Child’s Play (1988 - 2021)",
    "%subtitle5%": "Scream (1996 – 2022)",
    "%despedida%": "Hasta pronto . . . . ",
}

rectangle_replacement = {
    "IMG1": "image1.jpg",
    "IMG2": "image2.jpg",
    "IMG3": "image3.jpg",
    "IMG4": "image4.jpg",
    "IMG5": "image5.jpg",
}

def process_ppt(file, wildcard_text, rectangle_img):
    ppt = Presentation(file)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text in wildcard_text:
                replace_text(shape, wildcard_text)
            elif shape.shape_type == 1:
                replace_rect_img(slide, shape, rectangle_img)
 
    ppt.save('Resultado.pptx')

def replace_text(shape, wildcard):
        shape.text = wildcard[shape.text]

def replace_rect_img(slide, shape, rectangle):
    if shape.text in rectangle:
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height

        slide.shapes._spTree.remove(shape._element)
        slide.shapes.add_picture(rectangle[shape.text],left,top,width,height)

process_ppt(file_path, wildcard_replacement, rectangle_replacement)